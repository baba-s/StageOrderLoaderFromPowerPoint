from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PictureData:

	def __init__(self, filename, x, y):
		self.filename = filename
		self.x = x
		self.y = y

	def __str__(self):
		return self.filename + ", " + str(self.x) + ", " + str(self.y)


if __name__ == '__main__':

	ppt = Presentation("【.pptx のファイルパス】")

	pictureList = [[]]

	for slide in ppt.slides:
		pictureListPerSlide = []
		for shape in slide.shapes:
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				filename = shape._pic.nvPicPr.cNvPr.get('descr')
				x = shape.left
				y = shape.top
				data = PictureData(filename, x, y)
				pictureListPerSlide.append(data)
		pictureList.append(pictureListPerSlide)

	for pictureListPerSlide in pictureList:
		pictureListPerSlide.sort(key=lambda c: c.y)
		if len(pictureListPerSlide) == 10:

			# TODO:Python の API でキレイに書きたい
			splittedPictureListPerLines = [
				[pictureListPerSlide[0], pictureListPerSlide[1], pictureListPerSlide[2]],
				[pictureListPerSlide[3], pictureListPerSlide[4], pictureListPerSlide[5]],
				[pictureListPerSlide[6], pictureListPerSlide[7], pictureListPerSlide[8]],
				[pictureListPerSlide[9]],
			]
			for splittedPictureListPerLine in splittedPictureListPerLines:
				splittedPictureListPerLine.sort(key=lambda c: c.x)
				for pictureData in splittedPictureListPerLine:
					print(pictureData)
		print()
